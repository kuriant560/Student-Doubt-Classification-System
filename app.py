import streamlit as st
import os
import re
import html
from datetime import datetime
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from st_clipboard import copy_to_clipboard

import ollama

# Optional imports handled gracefully
try:
    import pdfplumber
    from docx import Document
    from pptx import Presentation
    EXTRACTION_READY = True
except ImportError:
    EXTRACTION_READY = False

# ================================
# PAGE CONFIGURATION
# ================================
st.set_page_config(
    page_title="Student Doubt Classification System",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ================================
# CORE LOGIC
# ================================
DATA_DIR = "data"
THRESHOLD = 0.15
OLLAMA_DOWNLOAD_URL = "https://ollama.com/download"


def inject_styles():
    st.markdown(
        """
        <style>
        :root {
            --bg: #07111f;
            --bg-soft: #101f33;
            --panel: rgba(10, 21, 36, 0.82);
            --panel-strong: rgba(16, 31, 51, 0.92);
            --border: rgba(143, 180, 255, 0.18);
            --text: #f5f7fb;
            --muted: #9eb2c9;
            --accent: #ff8f5c;
            --accent-2: #61d0ff;
            --warn: #ffd166;
        }

        .stApp {
            background:
                radial-gradient(circle at top left, rgba(97, 208, 255, 0.20), transparent 28%),
                radial-gradient(circle at top right, rgba(255, 143, 92, 0.20), transparent 26%),
                linear-gradient(180deg, #08111d 0%, #091523 55%, #07111f 100%);
            color: var(--text);
        }

        [data-testid="stSidebar"] {
            background:
                linear-gradient(180deg, rgba(13, 24, 40, 0.97) 0%, rgba(10, 18, 31, 0.97) 100%);
            border-right: 1px solid var(--border);
        }

        .block-container {
            padding-top: 2.2rem;
            padding-bottom: 3rem;
            max-width: 1180px;
        }

        h1, h2, h3 {
            font-family: "Avenir Next", "Trebuchet MS", sans-serif;
            letter-spacing: -0.03em;
        }

        .hero-shell {
            position: relative;
            overflow: hidden;
            padding: 2rem 2rem 1.8rem 2rem;
            border-radius: 28px;
            background:
                linear-gradient(140deg, rgba(255, 255, 255, 0.05), rgba(255, 255, 255, 0.02)),
                linear-gradient(135deg, rgba(10, 24, 42, 0.96), rgba(16, 31, 51, 0.88));
            border: 1px solid var(--border);
            box-shadow: 0 24px 70px rgba(0, 0, 0, 0.28);
            margin-bottom: 1.4rem;
        }

        .hero-shell::after {
            content: "";
            position: absolute;
            inset: auto -10% -40% auto;
            width: 260px;
            height: 260px;
            background: radial-gradient(circle, rgba(255, 143, 92, 0.24), transparent 64%);
            pointer-events: none;
        }

        .hero-kicker {
            display: inline-flex;
            align-items: center;
            gap: 0.45rem;
            font-size: 0.78rem;
            text-transform: uppercase;
            letter-spacing: 0.18em;
            color: var(--accent-2);
            margin-bottom: 0.9rem;
        }

        .hero-title {
            font-size: 3.05rem;
            line-height: 0.98;
            font-weight: 800;
            margin: 0;
            max-width: 8.5em;
        }

        .hero-copy {
            font-size: 1.03rem;
            line-height: 1.7;
            color: var(--muted);
            max-width: 44rem;
            margin: 0.9rem 0 1.15rem 0;
        }

        .hero-strip {
            display: flex;
            flex-wrap: wrap;
            gap: 0.7rem;
        }

        .hero-pill, .subject-pill {
            display: inline-flex;
            align-items: center;
            border: 1px solid rgba(143, 180, 255, 0.16);
            background: rgba(255, 255, 255, 0.04);
            border-radius: 999px;
            padding: 0.45rem 0.85rem;
            font-size: 0.86rem;
            color: var(--text);
        }

        .subject-pill {
            margin: 0 0.45rem 0.45rem 0;
            background: rgba(97, 208, 255, 0.08);
        }

        .input-shell, .result-shell {
            background: var(--panel);
            border: 1px solid var(--border);
            border-radius: 24px;
            padding: 1.15rem 1.2rem 1.25rem 1.2rem;
            box-shadow: 0 14px 44px rgba(0, 0, 0, 0.20);
        }

        .section-label {
            font-size: 0.78rem;
            text-transform: uppercase;
            letter-spacing: 0.16em;
            color: var(--accent);
            margin-bottom: 0.5rem;
        }

        .mini-note {
            color: var(--muted);
            font-size: 0.92rem;
            margin-bottom: 0.75rem;
        }

        .stat-card {
            height: 100%;
            padding: 1rem;
            border-radius: 18px;
            background: linear-gradient(180deg, rgba(255, 255, 255, 0.045), rgba(255, 255, 255, 0.02));
            border: 1px solid rgba(143, 180, 255, 0.14);
        }

        .stat-label {
            color: var(--muted);
            font-size: 0.78rem;
            text-transform: uppercase;
            letter-spacing: 0.14em;
            margin-bottom: 0.45rem;
        }

        .stat-value {
            color: var(--text);
            font-size: 1.05rem;
            font-weight: 700;
            line-height: 1.4;
        }

        .answer-card {
            border-radius: 22px;
            border: 1px solid rgba(143, 180, 255, 0.16);
            padding: 1.2rem 1.25rem;
            margin-top: 0.8rem;
            white-space: pre-wrap;
            line-height: 1.8;
        }

        .answer-card.retrieval {
            background: linear-gradient(180deg, rgba(97, 208, 255, 0.10), rgba(255, 255, 255, 0.03));
        }

        .answer-card.fallback {
            background: linear-gradient(180deg, rgba(255, 209, 102, 0.10), rgba(255, 255, 255, 0.03));
        }

        .notice-card {
            margin-top: 0.9rem;
            padding: 0.95rem 1rem;
            border-radius: 16px;
            background: rgba(255, 209, 102, 0.11);
            border: 1px solid rgba(255, 209, 102, 0.20);
            color: #ffe6a6;
        }

        .error-card {
            margin-top: 0.9rem;
            padding: 0.95rem 1rem;
            border-radius: 16px;
            background: rgba(255, 105, 97, 0.12);
            border: 1px solid rgba(255, 105, 97, 0.24);
            color: #ffd1ce;
        }

        .history-card {
            margin-top: 0.9rem;
            padding: 1rem 1.05rem;
            border-radius: 18px;
            background: linear-gradient(180deg, rgba(255, 255, 255, 0.045), rgba(255, 255, 255, 0.02));
            border: 1px solid rgba(143, 180, 255, 0.14);
        }

        .history-meta {
            color: var(--accent-2);
            font-size: 0.78rem;
            text-transform: uppercase;
            letter-spacing: 0.12em;
            margin-bottom: 0.55rem;
        }

        .history-question {
            font-weight: 700;
            color: var(--text);
            margin-bottom: 0.45rem;
        }

        .history-answer {
            color: var(--muted);
            line-height: 1.65;
        }

        .stTextArea textarea {
            min-height: 170px;
            border-radius: 18px;
            border: 1px solid rgba(143, 180, 255, 0.14);
            background: rgba(8, 16, 28, 0.86);
            color: var(--text);
            font-size: 1rem;
            box-shadow: inset 0 0 0 1px rgba(255, 255, 255, 0.02);
        }

        .stTextArea label, .stFileUploader label, .stButton button {
            font-family: "Avenir Next", "Trebuchet MS", sans-serif;
        }

        .stButton button {
            border: 0;
            border-radius: 999px;
            padding: 0.78rem 1.3rem;
            font-weight: 700;
            color: #08111d;
            background: linear-gradient(135deg, #ffb36b 0%, #ff7f67 100%);
            box-shadow: 0 12px 26px rgba(255, 127, 103, 0.32);
        }

        .stButton button:hover {
            color: #08111d;
            border: 0;
        }

        [data-testid="stExpander"] {
            border: 1px solid var(--border);
            border-radius: 18px;
            background: rgba(255, 255, 255, 0.03);
        }

        @media (max-width: 900px) {
            .hero-title {
                font-size: 2.35rem;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_hero(subject_count):
    st.markdown(
        f"""
        <section class="hero-shell">
            <div class="hero-kicker">Adaptive Study Assistant</div>
            <h1 class="hero-title">Student Doubt Classification System</h1>
            <p class="hero-copy">
                Search your notes, uploaded files, and built-in subject material through a retrieval-first workflow.
                When the match is weak, the interface can switch to local AI fallback through Ollama.
            </p>
            <div class="hero-strip">
                <span class="hero-pill">TF-IDF retrieval engine</span>
                <span class="hero-pill">{subject_count} subject sources loaded</span>
                <span class="hero-pill">TXT, PDF, DOCX, PPTX support</span>
            </div>
        </section>
        """,
        unsafe_allow_html=True,
    )


def render_stat_card(label, value):
    st.markdown(
        f"""
        <div class="stat-card">
            <div class="stat-label">{html.escape(label)}</div>
            <div class="stat-value">{html.escape(value)}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_answer_card(text, tone):
    escaped = html.escape(text).replace("\n", "<br>")
    st.markdown(
        f'<div class="answer-card {tone}">{escaped}</div>',
        unsafe_allow_html=True,
    )


def render_notice(message, tone="notice"):
    css_class = "error-card" if tone == "error" else "notice-card"
    st.markdown(
        f'<div class="{css_class}">{html.escape(message)}</div>',
        unsafe_allow_html=True,
    )


def render_history_item(entry):
    answer_preview = entry["answer"]
    if len(answer_preview) > 220:
        answer_preview = answer_preview[:220].rstrip() + "..."
    with st.sidebar.expander(entry["query"], expanded=False):
        st.caption(f'{entry["timestamp"]} • {entry["decision"]}')
        st.markdown(f'**Subject:** `{entry["subject"]}`')
        st.markdown("**Answer Preview**")
        st.write(answer_preview)

def chunk_text(text, chunk_size=500):
    """Character-based chunking."""
    chunks = []
    for i in range(0, len(text), chunk_size):
        chunk = text[i:i+chunk_size]
        if chunk.strip():
            chunks.append(chunk.strip())
    return chunks

def extract_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\\n"
    return text

def extract_docx(file):
    doc = Document(file)
    return "\\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def extract_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\\n"
    return text

def extract_txt(file):
    return file.read().decode("utf-8")

@st.cache_data
def process_documents(uploaded_files):
    corpus, sources, errors, success_count = [], [], [], 0
    if uploaded_files:
        for file in uploaded_files:
            try:
                file.seek(0)
                ext = file.name.split(".")[-1].lower()
                text = ""
                
                if ext == "txt": text = extract_txt(file)
                elif ext == "pdf": text = extract_pdf(file)
                elif ext == "docx": text = extract_docx(file)
                elif ext == "pptx": text = extract_pptx(file)
                else:
                    errors.append(f"File type not supported: {file.name}")
                    continue
                
                if not text.strip():
                    errors.append(f"No text extracted from PDF (empty content): {file.name}")
                    continue
                
                chunks = chunk_text(text)
                if chunks:
                    corpus.extend(chunks)
                    sources.extend([file.name] * len(chunks))
                    success_count += 1
                else:
                    errors.append(f"No valid chunks generated from {file.name}")
            except Exception as e:
                errors.append(f"Could not extract from {file.name}: {e}")
    else:
        if os.path.exists(DATA_DIR):
            for filename in os.listdir(DATA_DIR):
                if filename.endswith(".txt"):
                    path = os.path.join(DATA_DIR, filename)
                    with open(path, "r", encoding="utf-8") as f:
                        text = f.read()
                        chunks = chunk_text(text)
                        corpus.extend(chunks)
                        sources.extend([filename] * len(chunks))
    return corpus, sources, errors, success_count

def clean_output_text(text):
    text = text.replace("\\n", " ")
    text = re.sub(r'[^a-zA-Z0-9., ]', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def trim_text(text, max_len=400):
    return text[:max_len] + "..." if len(text) > max_len else text

def extract_key_sentences(text):
    sentences = text.split(".")
    return [s.strip() for s in sentences if s.strip()][:3]

def format_answer(text):
    sentences = extract_key_sentences(text)
    
    return f"""
Definition:
{sentences[0] + '.' if len(sentences) > 0 else ''}

Key Points:
- {sentences[1] + '.' if len(sentences) > 1 else ''}
- {sentences[2] + '.' if len(sentences) > 2 else ''}

Simple Explanation:
This concept explains how the system works in an easy and understandable way.
"""

def clean_text(text):
    text = str(text).lower()
    text = re.sub(r'[^a-zA-Z0-9 ]', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def expand_query(query):
    expansions = {
        "software design process": "software design process steps phases design methodology",
        "design process": "software design process phases steps",
    }
    for key in expansions:
        if key in query:
            return query + " " + expansions[key]
    return query


def get_best_match(query, corpus, sources):
    if not corpus: return None, 0.0, None, []
    
    processed_corpus = [clean_text(doc) for doc in corpus]
    expanded_query = expand_query(query)
    processed_query = clean_text(expanded_query)
    
    vectorizer = TfidfVectorizer(
        stop_words="english",
        ngram_range=(1,3),
        max_df=0.9,
        min_df=1
    )
    
    tfidf_matrix = vectorizer.fit_transform(processed_corpus + [processed_query])
    cosine_similarities = cosine_similarity(tfidf_matrix[-1], tfidf_matrix[:-1]).flatten()
    
    best_idx = cosine_similarities.argmax()
    best_score = cosine_similarities[best_idx]
    
    top_indices = cosine_similarities.argsort()[-3:][::-1]
    top_matches = [(corpus[i], cosine_similarities[i], sources[i]) for i in top_indices if cosine_similarities[i] > 0]
    
    return corpus[best_idx], best_score, sources[best_idx], top_matches

def generate_fallback(query):
    """Intelligent fallback utilizing local Ollama engine."""
    try:
        response = ollama.chat(
            model="llama3.2",
            messages=[
                {"role": "system", "content": "You are a helpful academic assistant. Keep your answers extremely simple, concise, and easy for a beginner student to understand. Avoid long walls of text and overly complex terminology."},
                {"role": "user", "content": query}
            ]
        )
        return response['message']['content']
    except Exception as e:
        return f"Error contacting local AI Fallback (Ollama): {str(e)}"


def ollama_response_failed(text):
    return isinstance(text, str) and text.startswith("Error contacting local AI Fallback (Ollama):")


def fallback_unavailable_message():
    return (
        "AI fallback is unavailable because Ollama is not running or not reachable on this machine. "
        f"Start Ollama locally and try again: {OLLAMA_DOWNLOAD_URL}"
    )

# Session State Init
if 'history' not in st.session_state: st.session_state.history = []
inject_styles()

# ================================
# SIDEBAR UI
# ================================
st.sidebar.title("📁 Document Upload")
uploaded_files = st.sidebar.file_uploader(
    "Upload Knowledge Base files", 
    type=["txt", "pdf", "docx", "pptx"], 
    accept_multiple_files=True
)

st.sidebar.markdown("### Uploaded Documents")
if uploaded_files:
    for f in uploaded_files:
        st.sidebar.markdown(
            f'<div class="subject-pill">{html.escape(f.name)}</div>',
            unsafe_allow_html=True,
        )
else:
    st.sidebar.info("Using default files from data/ folder.")

corpus, sources, errors, success_count = process_documents(uploaded_files)

if uploaded_files:
    if success_count > 0:
        st.sidebar.success(f"Successfully processed {success_count} files.")
    for err in errors:
        st.sidebar.warning(err)
    if not EXTRACTION_READY:
        st.sidebar.error("Warning: Advanced parsers missing. Run `pip install pdfplumber python-docx python-pptx`")


st.sidebar.markdown("### Detected Subjects")
unique_subjects = sorted(list(set([os.path.splitext(s)[0].replace("_", " ") for s in sources])))
if unique_subjects:
    for sub in unique_subjects:
        st.sidebar.markdown(
            f'<div class="subject-pill">{html.escape(sub)}</div>',
            unsafe_allow_html=True,
        )
else:
    st.sidebar.markdown('<div class="subject-pill">General</div>', unsafe_allow_html=True)

st.sidebar.markdown("### Recent Searches")
if st.session_state.history:
    for entry in st.session_state.history:
        render_history_item(entry)
else:
    st.sidebar.info("Your recent searches will appear here.")

# ================================
# MAIN SCREEN UI
# ================================
render_hero(len(unique_subjects) if unique_subjects else 1)
st.markdown('<div class="section-label">Question Workspace</div>', unsafe_allow_html=True)
st.markdown(
    '<div class="mini-note">Ask a concept question, process uploaded notes, and let the app decide whether retrieval or local AI should answer it.</div>',
    unsafe_allow_html=True,
)
st.divider()

# Input
query = st.text_area("Ask your academic doubt...", placeholder="Enter question here...")

if st.button("Get Answer", type="primary"):
    if not query.strip():
        st.warning("Please enter a question.")
    else:
        with st.spinner("Processing using NLP..."):
            
            # Step 5 Validation
            if len(corpus) == 0:
                st.error("No valid documents extracted. Cannot compute similarity. Please upload better documents.")
                st.stop()
                
            best_doc, confidence, source_file, top_matches = get_best_match(query, corpus, sources)
            
            is_fallback = confidence < THRESHOLD
            
            # Hybrid logic branches
            if is_fallback:
                subject = "GENERAL"
                with st.spinner("Generating AI response..."):
                    decision_label = "Switched to AI Fallback"
                    source_tag = "🤖 AI Generated Answer"
                    answer_text = generate_fallback(query)
            else:
                subject = os.path.splitext(source_file)[0].replace("_", " ").upper()
                decision_label = "Matched via Document Retrieval"
                source_tag = "📄 Answer from Documents"
                
                with st.spinner("Preparing answer from matched document..."):
                    cleaned = clean_output_text(best_doc)
                    trimmed = trim_text(cleaned)
                    
                    prompt = f"""
Convert the following academic content into a structured student answer with:
1. Definition
2. 3 bullet points
3. Simple explanation

Text:
{trimmed}
"""
                    try:
                        ai_rewritten_answer = generate_fallback(prompt)
                        if ollama_response_failed(ai_rewritten_answer):
                            answer_text = format_answer(trimmed)
                        else:
                            answer_text = ai_rewritten_answer
                    except Exception:
                        answer_text = format_answer(trimmed)

        st.session_state.history.insert(
            0,
            {
                "query": query,
                "answer": answer_text,
                "timestamp": datetime.now().strftime("%I:%M %p"),
                "decision": decision_label,
                "subject": subject,
            },
        )
        st.session_state.history = st.session_state.history[:5]

        # --------------------------
        # UI Rendering
        # --------------------------
        st.markdown('<div class="section-label">Results</div>', unsafe_allow_html=True)
        st.markdown(
            f'<div class="mini-note"><strong>Your Question:</strong> {html.escape(query)}</div>',
            unsafe_allow_html=True,
        )

        metric_cols = st.columns(3)
        with metric_cols[0]:
            render_stat_card("Subject", subject)
        with metric_cols[1]:
            render_stat_card(
                "Confidence",
                f"{round(confidence * 100, 2)}%" if not is_fallback else "AI Fallback",
            )
        with metric_cols[2]:
            render_stat_card("Decision", decision_label)
        
        # DEBUG OUTPUT
        with st.expander("🛠️ Advanced Debug Info"):
            st.info(f"""
            **NLP Pipeline Debug Info:**
            - Total Chunks Mapped: `{len(corpus)}`
            - Highest Similarity Score: `{confidence:.4f}`
            - Similarity Threshold: `{THRESHOLD:.2f}`
            """)
            
            if top_matches:
                for idx, (match_doc, match_score, match_source) in enumerate(top_matches):
                    st.markdown(f"**Match {idx+1}** | Score: `{match_score:.4f}` | Source: `{match_source}`\\n\\n> {match_doc[:300]}...")
            elif len(corpus) > 0:
                st.write(corpus[0] if len(corpus) > 0 else "Empty")
        
        if confidence > 0:
            st.progress(float(confidence))

        # Decision Output
        st.markdown("---")
        st.markdown(
            f'<div class="mini-note"><strong>Source Tag:</strong> {html.escape(source_tag)}</div>',
            unsafe_allow_html=True,
        )
        st.markdown("### Answer")
        
        # Visual color cue for the answer output
        if is_fallback:
            if ollama_response_failed(answer_text):
                render_notice(fallback_unavailable_message(), tone="error")
            else:
                render_answer_card(answer_text, "fallback")
            render_notice("Low confidence answer. Try rephrasing your query or upload better documents.")
        else:
            render_answer_card(answer_text, "retrieval")
            
        # Try clipboard standard
        try:
            copy_to_clipboard(answer_text, text="📋 Copy", key=f"copy_{len(st.session_state.history)}")
        except Exception:
            pass
